from __future__ import annotations

import html
import json
import math
import re
import shutil
import subprocess
import sys
import tempfile
from collections import Counter, defaultdict
from concurrent.futures import ThreadPoolExecutor
from datetime import date, datetime, timedelta, timezone
from pathlib import Path
from statistics import mean, median
from urllib.request import Request, urlopen

import fitz
from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from PIL import Image, ImageDraw, ImageStat
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")


ROOT = Path(__file__).resolve().parent
OUTPUT_DIR = ROOT / "output" / "director_manager_report"
TO_MANAGER_DIR = OUTPUT_DIR / "to_send_governorates_simple"
DATA_DIR = OUTPUT_DIR / "data"
QA_DIR = OUTPUT_DIR / "qa"
ASSETS_DIR = ROOT / "output" / "assets"

LOGO_PATH = ASSETS_DIR / "alforaij_logo.png"
COVER_BG_PATH = ASSETS_DIR / "kuwait_glass_cover.png"

REPORT_HTML_PATH = OUTPUT_DIR / "director_manager_report_source.html"
REPORT_PDF_PATH = OUTPUT_DIR / "director_manager_report.pdf"
PPTX_PATH = OUTPUT_DIR / "director_manager_report.pptx"
EXCEL_PATH = OUTPUT_DIR / "director_manager_data_register.xlsx"
SOURCES_HTML_PATH = OUTPUT_DIR / "sources_methodology.html"
SOURCES_PDF_PATH = OUTPUT_DIR / "sources_methodology.pdf"
QA_MD_PATH = OUTPUT_DIR / "QA.md"
MANIFEST_PATH = OUTPUT_DIR / "manifest.json"
MONTAGE_PATH = QA_DIR / "pdf_pages_montage.png"

FINAL_REPORT_PDF = TO_MANAGER_DIR / "01_تقرير_إحصائيات_الفريج_حسب_المحافظات.pdf"
FINAL_PPTX = TO_MANAGER_DIR / "02_عرض_إحصائيات_الفريج_حسب_المحافظات_قابل_للتعديل.pptx"
FINAL_EXCEL = TO_MANAGER_DIR / "03_سجل_الأرقام_والمصادر_الفريج.xlsx"
FINAL_SOURCES_PDF = TO_MANAGER_DIR / "04_مصادر_ومنهجية_البيانات.pdf"

SNAPSHOT_DATE = date(2026, 7, 29)
SNAPSHOT_LABEL = "حتى 29 يوليو 2026"
CLIENT_NAME = "شركة عبدالعزيز سعود الفريج العقارية"
REPORT_TITLE = "إحصائيات الفريج التشغيلية"
REPORT_SUBTITLE = "حركة الدلال | عروض البيع والشراء | البيوت للبيع | الإيجارات"
SOURCE_NAME = "منصة الفريج - واجهة البحث العامة"

ALFORAIJ_FRONT_URL = "https://front.alforaij.com/"
ALFORAIJ_SEARCH_URL = "https://search.alforaij.com/"
TRANSACTION_TYPES_URL = "https://search.alforaij.com/api/transactiontypes"
LISTINGS_URL = "https://search.alforaij.com/api/internallistings/search"
DETAIL_URL = "https://search.alforaij.com/api/internallistings"
PUBLIC_DETAIL_URL = "https://front.alforaij.com/Listing/Detail"

BRAND = {
    "dark": "111827",
    "dark_2": "1F2937",
    "ink": "0F172A",
    "muted": "475569",
    "soft": "F8FAFC",
    "line": "CBD5E1",
    "gold": "D97706",
    "gold_2": "F59E0B",
    "blue": "1D4ED8",
    "green": "047857",
    "red": "B91C1C",
    "teal": "0F766E",
    "white": "FFFFFF",
}

FONT_AR = "Arial"
CORE_TRANSACTION_IDS = [1, 2, 3, 4]
ALL_TRANSACTION_IDS = [1, 2, 3, 4, 5]
HOUSE_TERMS = ("بيت", "فيلا", "ڤيلا", "قسيمة", "منزل")
FEATURE_TERMS = (
    "زاوية",
    "شارع وسكة",
    "سكة",
    "شارعين",
    "شارع",
    "رأس",
    "بطن وظهر",
    "ارتداد",
    "ساحة",
    "قرب الخدمات",
    "مدخل ومخرج",
)


def ensure_dirs() -> None:
    for path in [OUTPUT_DIR, TO_MANAGER_DIR, DATA_DIR, QA_DIR]:
        path.mkdir(parents=True, exist_ok=True)


def hex_to_rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def css_hex(hex_value: str) -> str:
    return "#" + hex_value.lstrip("#")


def fmt_int(value: int | float | None) -> str:
    if value is None:
        return "0"
    return f"{int(round(float(value))):,}"


def fmt_decimal(value: int | float | None, decimals: int = 1) -> str:
    if value is None:
        return "0"
    return f"{float(value):,.{decimals}f}"


def fmt_kd(value: int | float | None) -> str:
    if value is None:
        return "غير معلن"
    rounded = int(round(float(value)))
    if abs(rounded) >= 10_000:
        return f"{int(round(rounded / 1000))} ألف"
    return str(rounded)


def pct(part: int | float, whole: int | float) -> float:
    if not whole:
        return 0.0
    return round((float(part) / float(whole)) * 100, 1)


def fetch_json(url: str, timeout: int = 60) -> dict | list:
    req = Request(
        url,
        headers={
            "User-Agent": "Mozilla/5.0",
            "Accept": "application/json,text/plain,*/*",
        },
    )
    with urlopen(req, timeout=timeout) as response:
        return json.loads(response.read().decode("utf-8"))


def parse_datetime(value: str | None) -> datetime | None:
    if not value:
        return None
    text = str(value).strip()
    if not text:
        return None
    if "." in text:
        head, tail = text.split(".", 1)
        match = re.match(r"(\d+)(.*)", tail)
        if match:
            text = f"{head}.{match.group(1)[:6]}{match.group(2)}"
    try:
        return datetime.fromisoformat(text)
    except ValueError:
        return None


def clean_text(value: object) -> str:
    text = str(value or "")
    text = re.sub(r"[\r\n\t]+", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def strip_private_numbers(text: str) -> str:
    text = re.sub(r"\+?\d[\d\s\-]{5,}\d", "", text)
    return re.sub(r"\s+", " ", text).strip()


def strip_contact_phrases(text: str) -> str:
    text = re.split(r"(?:للاستفسار|للتواصل|واتساب|whatsapp|alforaij\.com|📞|📱|🌐)", text, maxsplit=1, flags=re.IGNORECASE)[0]
    return strip_private_numbers(clean_text(text))


def get_field(item: dict, *names: str, default: object = "") -> object:
    for name in names:
        if name in item and item[name] not in (None, ""):
            return item[name]
    return default


def parse_price(value: object) -> float:
    try:
        return float(value or 0)
    except (TypeError, ValueError):
        return 0.0


def extract_area_from_text(text: str) -> float | None:
    patterns = [
        r"(?:المساحة|مساحة|مساحه|مساحة العقار|مساحة البيت|مساحته|مساحتها)\s*[:：]?\s*([0-9٠-٩,\.]+)\s*(?:متر مربع|م²|متر|مترمربع)?",
    ]
    translation = str.maketrans("٠١٢٣٤٥٦٧٨٩", "0123456789")
    for pattern in patterns:
        match = re.search(pattern, text)
        if not match:
            continue
        value = match.group(1).translate(translation).replace(",", "")
        try:
            area = float(value)
        except ValueError:
            continue
        if 10 <= area <= 100000:
            return area
    return None


def extract_after_label(text: str, labels: tuple[str, ...], max_len: int = 80) -> str:
    for label in labels:
        match = re.search(rf"{label}\s*[:：]\s*([^\n\r▪]+)", text)
        if match:
            return strip_private_numbers(clean_text(match.group(1)))[:max_len]
    return ""


def extract_listing_mode(text: str) -> str:
    mode = extract_after_label(text, ("نوع الإعلان",), 40)
    if mode:
        return mode
    if "مباشر" in text:
        return "مباشر"
    return "غير محدد"


def classify_detail(text: str, property_type_name: str, transaction_type_id: int) -> str:
    combined = f"{property_type_name} {text}"
    if transaction_type_id in (3, 4):
        if any(term in combined for term in HOUSE_TERMS):
            return "طلب بيت/فيلا"
        return "طلب عقاري"
    if "شقة" in combined:
        return "شقة"
    if "دور" in combined:
        return "دور"
    if any(term in combined for term in ("محل", "مكتب", "تجاري")):
        return "تجاري"
    if any(term in combined for term in HOUSE_TERMS):
        return "بيت/فيلا/قسيمة"
    return property_type_name or "غير محدد"


def extract_features(text: str) -> str:
    found = []
    for term in FEATURE_TERMS:
        if term in text and term not in found:
            found.append(term)
    return "، ".join(found[:5])


def normalize_listing(item: dict, type_id: int) -> dict:
    published_dt = parse_datetime(get_field(item, "published_at", "publishedAt", "created_at", "createdAt", default=None))
    published_date = published_dt.date() if published_dt else None
    days_since = (SNAPSHOT_DATE - published_date).days if published_date else None
    if days_since is not None and days_since < 0:
        days_since = 0

    title_raw = str(get_field(item, "title_ar", "titleAr", "title", default="") or "")
    description_raw = str(get_field(item, "description_ar", "descriptionAr", "description", default="") or "")
    private_detail_raw = str(get_field(item, "ownerName", "owner_name", default="") or "")
    detail_text_raw = "\n".join([title_raw, description_raw, private_detail_raw])
    title = strip_contact_phrases(title_raw)
    description = strip_contact_phrases(description_raw)
    price = parse_price(get_field(item, "price", default=0))
    area = parse_price(get_field(item, "square_meters", "squareMeters", "area", default=0))
    extracted_area = area if area > 0 else extract_area_from_text(detail_text_raw)
    property_type_name = clean_text(get_field(item, "property_type_name", "propertyTypeName", default="غير محدد"))
    detail_class = classify_detail(detail_text_raw, property_type_name, type_id)
    features = extract_features(detail_text_raw)
    condition = strip_contact_phrases(extract_after_label(detail_text_raw, ("الحالة", "المواصفات", "الوصف", "الموقع والمميزات"), 90))
    listing_mode = extract_listing_mode(detail_text_raw)

    item_id = int(get_field(item, "id", default=0) or 0)
    detail_loaded = bool(get_field(item, "_detail_loaded", default=False))
    detail_url = f"{DETAIL_URL}/{item_id}"
    public_detail_url = f"{PUBLIC_DETAIL_URL}/{item_id}"
    search_url = f"{LISTINGS_URL}?page=1&pageSize=100&transactionType={type_id}"

    return {
        "id": item_id,
        "code": f"AF-{item_id}",
        "transaction_type_id": type_id,
        "transaction_type_name": clean_text(get_field(item, "transaction_type_name", "transactionTypeName", default="")),
        "property_type_name": property_type_name,
        "detail_class": detail_class,
        "governorate_name": clean_text(get_field(item, "governorate_name", "governorateName", default="غير محدد")),
        "city_name": clean_text(get_field(item, "city_name", "cityName", default="غير محدد")),
        "price": price,
        "is_priced": price > 0,
        "area": extracted_area if extracted_area and extracted_area > 0 else None,
        "published_date": published_date,
        "published_at": published_dt.isoformat() if published_dt else "",
        "days_since": days_since,
        "listing_mode": listing_mode,
        "condition": condition,
        "features": features,
        "has_image": bool(clean_text(get_field(item, "featuredImagePath", "featured_image_path", default="")) or clean_text(get_field(item, "imagesJson", "images_json", default=""))),
        "detail_loaded": detail_loaded,
        "title_clean": title[:160],
        "description_clean": description[:200],
        "detail_url": detail_url,
        "source_url": public_detail_url if item_id else search_url,
    }


def fetch_listings_for_type(type_id: int) -> list[dict]:
    page_size = 100
    first_url = f"{LISTINGS_URL}?page=1&pageSize={page_size}&transactionType={type_id}"
    first = fetch_json(first_url)
    data = list(first.get("data", []))
    meta = first.get("meta", {}) or {}
    last_page = int(meta.get("last_page", 1) or 1)
    for page in range(2, last_page + 1):
        url = f"{LISTINGS_URL}?page={page}&pageSize={page_size}&transactionType={type_id}"
        payload = fetch_json(url)
        data.extend(payload.get("data", []))

    def with_detail(item: dict) -> dict:
        item_id = int(get_field(item, "id", default=0) or 0)
        merged = dict(item)
        if item_id:
            try:
                detail_payload = fetch_json(f"{DETAIL_URL}/{item_id}", timeout=18)
                detail = detail_payload.get("data", detail_payload) if isinstance(detail_payload, dict) else {}
                if isinstance(detail, dict):
                    merged.update(detail)
                    merged["_detail_loaded"] = True
            except Exception:
                merged["_detail_loaded"] = False
        return merged

    with ThreadPoolExecutor(max_workers=16) as executor:
        detailed_data = list(executor.map(with_detail, data))

    seen: set[int] = set()
    normalized: list[dict] = []
    for item in detailed_data:
        listing = normalize_listing(item, type_id)
        if listing["id"] in seen:
            continue
        seen.add(listing["id"])
        normalized.append(listing)
    normalized.sort(key=lambda row: (row["published_date"] or date.min, row["id"]), reverse=True)
    return normalized


def fetch_all_data() -> tuple[dict[int, str], list[dict]]:
    raw_types = fetch_json(TRANSACTION_TYPES_URL)
    type_names: dict[int, str] = {}
    for item in raw_types:
        type_id = int(get_field(item, "id", default=0) or 0)
        type_names[type_id] = clean_text(get_field(item, "name_ar", "nameAr", "title_ar", "titleAr", "name", default=str(type_id)))
    listings: list[dict] = []
    for type_id in ALL_TRANSACTION_IDS:
        listings.extend(fetch_listings_for_type(type_id))
    for row in listings:
        if not row["transaction_type_name"]:
            row["transaction_type_name"] = type_names.get(row["transaction_type_id"], str(row["transaction_type_id"]))
    return type_names, listings


def is_house_sale(row: dict) -> bool:
    if row["transaction_type_id"] != 1:
        return False
    text = " ".join(
        [
            row.get("property_type_name", ""),
            row.get("title_clean", ""),
            row.get("description_clean", ""),
        ]
    )
    return any(term in text for term in HOUSE_TERMS)


def price_stats(rows: list[dict]) -> dict:
    values = [row["price"] for row in rows if row.get("is_priced")]
    return {
        "priced_count": len(values),
        "undisclosed_count": len(rows) - len(values),
        "avg": round(mean(values), 1) if values else None,
        "median": round(median(values), 1) if values else None,
        "min": round(min(values), 1) if values else None,
        "max": round(max(values), 1) if values else None,
    }


def top_counts(rows: list[dict], field: str, limit: int = 5) -> list[dict]:
    counter = Counter(row.get(field) or "غير محدد" for row in rows)
    return [{"label": label, "value": value} for label, value in counter.most_common(limit)]


def type_counts(rows: list[dict], type_names: dict[int, str], include_exchange: bool = False) -> list[dict]:
    ids = ALL_TRANSACTION_IDS if include_exchange else CORE_TRANSACTION_IDS
    counter = Counter(row["transaction_type_id"] for row in rows)
    return [{"label": type_names.get(type_id, str(type_id)), "value": counter.get(type_id, 0)} for type_id in ids]


def governorate_summary(core_rows: list[dict], house_sale_rows: list[dict]) -> list[dict]:
    governorates = sorted({row.get("governorate_name") or "غير محدد" for row in core_rows})
    house_counts = Counter(row.get("governorate_name") or "غير محدد" for row in house_sale_rows)
    result = []
    for governorate in governorates:
        rows = [row for row in core_rows if (row.get("governorate_name") or "غير محدد") == governorate]
        sale = sum(1 for row in rows if row["transaction_type_id"] == 1)
        rent = sum(1 for row in rows if row["transaction_type_id"] == 2)
        buy = sum(1 for row in rows if row["transaction_type_id"] == 3)
        rent_request = sum(1 for row in rows if row["transaction_type_id"] == 4)
        sale_gap = sale - buy
        rent_gap = rent - rent_request
        result.append(
            {
                "label": governorate,
                "total": len(rows),
                "sale": sale,
                "buy": buy,
                "sale_gap": sale_gap,
                "sale_ratio": round(sale / buy, 1) if buy else None,
                "house_sale": house_counts.get(governorate, 0),
                "rent": rent,
                "rent_request": rent_request,
                "rent_gap": rent_gap,
                "rent_ratio": round(rent / rent_request, 1) if rent_request else None,
            }
        )
    result.sort(key=lambda row: row["total"], reverse=True)
    return result


def price_bands(rows: list[dict], mode: str) -> list[dict]:
    if mode == "rent":
        bands = [
            ("أقل من 250", 1, 249),
            ("250-399", 250, 399),
            ("400-549", 400, 549),
            ("550-799", 550, 799),
            ("800 فأكثر", 800, math.inf),
        ]
    else:
        bands = [
            ("أقل من 150 ألف", 1, 149999),
            ("150-249 ألف", 150000, 249999),
            ("250-349 ألف", 250000, 349999),
            ("350-499 ألف", 350000, 499999),
            ("500 ألف فأكثر", 500000, math.inf),
        ]
    result = []
    for label, low, high in bands:
        count = sum(1 for row in rows if row.get("is_priced") and low <= row["price"] <= high)
        result.append({"label": label, "value": count})
    undisclosed = sum(1 for row in rows if not row.get("is_priced"))
    if undisclosed:
        result.append({"label": "غير معلن", "value": undisclosed})
    return result


def recent_count(rows: list[dict], days: int) -> int:
    cutoff = SNAPSHOT_DATE - timedelta(days=days)
    return sum(1 for row in rows if row.get("published_date") and row["published_date"] >= cutoff)


def build_analysis(type_names: dict[int, str], listings: list[dict]) -> dict:
    by_type = {type_id: [row for row in listings if row["transaction_type_id"] == type_id] for type_id in ALL_TRANSACTION_IDS}
    core_rows = [row for row in listings if row["transaction_type_id"] in CORE_TRANSACTION_IDS]
    sale_rows = by_type[1]
    rent_rows = by_type[2]
    buy_rows = by_type[3]
    rent_request_rows = by_type[4]
    exchange_rows = by_type[5]
    house_sale_rows = [row for row in sale_rows if is_house_sale(row)]

    sale_stats = price_stats(sale_rows)
    buy_stats = price_stats(buy_rows)
    house_stats = price_stats(house_sale_rows)
    rent_stats = price_stats(rent_rows)
    rent_req_stats = price_stats(rent_request_rows)
    gov_summary = governorate_summary(core_rows, house_sale_rows)

    total_core = len(core_rows)
    recent_7 = recent_count(core_rows, 7)
    recent_30 = recent_count(core_rows, 30)
    details_loaded = sum(1 for row in listings if row.get("detail_loaded"))
    direct_count = sum(1 for row in core_rows if "مباشر" in row.get("listing_mode", ""))

    sale_to_buy_ratio = round(len(sale_rows) / len(buy_rows), 1) if buy_rows else None
    rent_offer_to_request_ratio = round(len(rent_rows) / len(rent_request_rows), 1) if rent_request_rows else None

    monthly_projection = {
        "إجمالي الحركة": recent_30,
        "للبيع": recent_count(sale_rows, 30),
        "للإيجار": recent_count(rent_rows, 30),
        "مطلوب للشراء": recent_count(buy_rows, 30),
        "مطلوب للإيجار": recent_count(rent_request_rows, 30),
    }

    metrics = {
        "total_core": total_core,
        "total_all": len(listings),
        "sale_count": len(sale_rows),
        "rent_count": len(rent_rows),
        "buy_count": len(buy_rows),
        "rent_request_count": len(rent_request_rows),
        "exchange_count": len(exchange_rows),
        "house_sale_count": len(house_sale_rows),
        "details_loaded": details_loaded,
        "direct_count": direct_count,
        "recent_7": recent_7,
        "recent_30": recent_30,
        "sale_to_buy_ratio": sale_to_buy_ratio,
        "rent_offer_to_request_ratio": rent_offer_to_request_ratio,
        "sale_stats": sale_stats,
        "buy_stats": buy_stats,
        "house_stats": house_stats,
        "rent_stats": rent_stats,
        "rent_request_stats": rent_req_stats,
        "type_counts": type_counts(listings, type_names),
        "type_counts_all": type_counts(listings, type_names, include_exchange=True),
        "movement_city_top": top_counts(core_rows, "city_name", 7),
        "movement_class_top": top_counts(core_rows, "detail_class", 7),
        "listing_mode_top": top_counts(core_rows, "listing_mode", 5),
        "sale_city_top": top_counts(sale_rows, "city_name", 6),
        "buy_city_top": top_counts(buy_rows, "city_name", 6),
        "house_city_top": top_counts(house_sale_rows, "city_name", 6),
        "rent_city_top": top_counts(rent_rows, "city_name", 6),
        "house_class_top": top_counts(house_sale_rows, "detail_class", 6),
        "rent_class_top": top_counts(rent_rows, "detail_class", 6),
        "sale_governorate_top": top_counts(sale_rows, "governorate_name", 6),
        "rent_governorate_top": top_counts(rent_rows, "governorate_name", 6),
        "governorate_summary": gov_summary,
        "governorate_movement_top": [{"label": row["label"], "value": row["total"]} for row in gov_summary],
        "governorate_sale_top": [{"label": row["label"], "value": row["sale"]} for row in gov_summary],
        "governorate_buy_top": [{"label": row["label"], "value": row["buy"]} for row in gov_summary],
        "governorate_house_top": [{"label": row["label"], "value": row["house_sale"]} for row in gov_summary],
        "governorate_rent_top": [{"label": row["label"], "value": row["rent"]} for row in gov_summary],
        "governorate_rent_request_top": [{"label": row["label"], "value": row["rent_request"]} for row in gov_summary],
        "governorate_rent_gap_top": sorted(
            [{"label": row["label"], "value": row["rent_gap"]} for row in gov_summary],
            key=lambda row: row["value"],
            reverse=True,
        ),
        "house_price_bands": price_bands(house_sale_rows, "sale"),
        "rent_price_bands": price_bands(rent_rows, "rent"),
        "monthly_projection": monthly_projection,
        "by_type": by_type,
        "core_rows": core_rows,
        "sale_rows": sale_rows,
        "buy_rows": buy_rows,
        "rent_rows": rent_rows,
        "rent_request_rows": rent_request_rows,
        "exchange_rows": exchange_rows,
        "house_sale_rows": house_sale_rows,
        "type_names": type_names,
    }
    return metrics


def bar_chart_html(data: list[dict], suffix: str = "", max_value: float | None = None, color: str | None = None) -> str:
    if not data:
        return '<div class="empty">لا توجد قيم قابلة للعرض في هذا المحور.</div>'
    max_val = max_value or max(item["value"] for item in data) or 1
    accent = color or css_hex(BRAND["blue"])
    rows = []
    for item in data:
        width = max(3, (float(item["value"]) / max_val) * 100)
        rows.append(
            f"""
            <div class="bar-row">
              <span class="bar-label">{html.escape(str(item["label"]))}</span>
              <span class="bar-track"><i style="width:{width:.1f}%;background:{accent}"></i></span>
              <strong>{fmt_int(item["value"])} {html.escape(suffix)}</strong>
            </div>
            """
        )
    return '<div class="bars">' + "\n".join(rows) + "</div>"


def band_cards_html(data: list[dict], suffix: str = "") -> str:
    cards = []
    for item in data:
        cards.append(
            f"""
            <div class="band-card">
              <span>{html.escape(str(item["label"]))}</span>
              <strong>{fmt_int(item["value"])}</strong>
            </div>
            """
        )
    return '<div class="band-grid">' + "\n".join(cards) + "</div>"


def governorate_table_html(rows: list[dict]) -> str:
    headers = ["المحافظة", "إجمالي الحركة", "للبيع", "شراء", "بيوت للبيع", "إيجار", "طلب إيجار"]
    body_rows = [
        [
            row["label"],
            fmt_int(row["total"]),
            fmt_int(row["sale"]),
            fmt_int(row["buy"]),
            fmt_int(row["house_sale"]),
            fmt_int(row["rent"]),
            fmt_int(row["rent_request"]),
        ]
        for row in rows
    ]
    header_html = "".join(f"<th>{html.escape(header)}</th>" for header in headers)
    row_html = []
    for row in body_rows:
        row_html.append("<tr>" + "".join(f"<td>{html.escape(str(cell))}</td>" for cell in row) + "</tr>")
    return f"""<table class="gov-table"><thead><tr>{header_html}</tr></thead><tbody>{"".join(row_html)}</tbody></table>"""


def kpi_html(title: str, value: str, unit: str, note: str = "", accent: str | None = None) -> str:
    style = f' style="border-right-color:{accent}"' if accent else ""
    note_html = f"<small>{html.escape(note)}</small>" if note else ""
    return f"""
    <article class="kpi"{style}>
      <span>{html.escape(title)}</span>
      <strong>{html.escape(value)}</strong>
      <em>{html.escape(unit)}</em>
      {note_html}
    </article>
    """


def source_footer(page: int) -> str:
    return (
        f"<footer><span>{html.escape(CLIENT_NAME)} | المصدر: {html.escape(SOURCE_NAME)}</span>"
        f"<span>{page:02d}</span></footer>"
    )


def slide_header(title: str, subtitle: str) -> str:
    logo = "../assets/alforaij_logo.png" if LOGO_PATH.exists() else ""
    logo_html = f'<img class="logo" src="{logo}" alt="{html.escape(CLIENT_NAME)}">' if logo else ""
    return f"""
    <header>
      <div class="title-block">
        <h1>{html.escape(title)}</h1>
        <p>{html.escape(subtitle)}</p>
      </div>
      {logo_html}
    </header>
    """


def build_report_html(metrics: dict) -> None:
    sale_stats = metrics["sale_stats"]
    buy_stats = metrics["buy_stats"]
    house_stats = metrics["house_stats"]
    rent_stats = metrics["rent_stats"]
    ratio = metrics["sale_to_buy_ratio"]
    rent_ratio = metrics["rent_offer_to_request_ratio"]

    sections = []
    sections.append(
        f"""
        <section class="slide cover">
          <img class="cover-logo" src="../assets/alforaij_logo.png" alt="{html.escape(CLIENT_NAME)}">
          <div class="cover-glass">
            <p class="cover-kicker">{html.escape(SNAPSHOT_LABEL)}</p>
            <h1>{html.escape(REPORT_TITLE)}</h1>
            <h2>{html.escape(REPORT_SUBTITLE)}</h2>
            <p>قراءة تشغيلية مختصرة من السجلات المنشورة على منصة الفريج، مع إخفاء بيانات التواصل وأي بيانات شخصية.</p>
            <div class="cover-chips">
              <span><b>{fmt_int(metrics["total_core"])}</b>حركة ضمن نطاق الطلب</span>
              <span><b>{fmt_int(metrics["sale_count"])}</b>عرض بيع</span>
              <span><b>{fmt_int(metrics["rent_count"])}</b>عرض إيجار</span>
              <span><b>{fmt_int(metrics["buy_count"])}</b>طلب شراء</span>
            </div>
          </div>
          {source_footer(1)}
        </section>
        """
    )

    sections.append(
        f"""
        <section class="slide">
          {slide_header("الملخص التنفيذي", "الأرقام التي تجيب مباشرة عن المطلوب")}
          <main>
            <div class="kpi-grid six">
              {kpi_html("حركة الدلال", fmt_int(metrics["total_core"]), "إعلان/طلب منشور", f"{fmt_int(metrics['details_loaded'])} إعلان بتفاصيل متاحة")}
              {kpi_html("عروض البيع", fmt_int(metrics["sale_count"]), "عرض", f"{fmt_int(sale_stats['priced_count'])} بسعر معلن", css_hex(BRAND["blue"]))}
              {kpi_html("طلبات الشراء", fmt_int(metrics["buy_count"]), "طلب", f"مقابل كل طلب شراء: {fmt_decimal(ratio, 1)} عرض بيع", css_hex(BRAND["green"]))}
              {kpi_html("بيوت للبيع", fmt_int(metrics["house_sale_count"]), "بيت/فيلا/قسيمة", f"متوسط معلن {fmt_kd(house_stats['avg'])} د.ك", css_hex(BRAND["gold"]))}
              {kpi_html("عروض الإيجار", fmt_int(metrics["rent_count"]), "عرض", f"متوسط معلن {fmt_kd(rent_stats['avg'])} د.ك شهرياً", css_hex(BRAND["teal"]))}
              {kpi_html("طلبات الإيجار", fmt_int(metrics["rent_request_count"]), "طلب", f"نسبة العرض للطلب {fmt_decimal(rent_ratio, 1)} مرة", css_hex(BRAND["red"]))}
            </div>
          </main>
          {source_footer(2)}
        </section>
        """
    )

    sections.append(
        f"""
        <section class="slide">
          {slide_header("حركة الدلال", "قياس النشاط المنشور على المنصة")}
          <main>
            <div class="split-panels">
              <div class="panel">
                <h2>الحركة حسب نوع المعاملة</h2>
                {bar_chart_html(metrics["type_counts"], "سجل", color=css_hex(BRAND["gold"]))}
              </div>
              <div class="panel">
                <h2>أعلى مناطق الحركة</h2>
                {bar_chart_html(metrics["movement_city_top"], "سجل", color=css_hex(BRAND["blue"]))}
              </div>
            </div>
          </main>
          {source_footer(3)}
        </section>
        """
    )

    sections.append(
        f"""
        <section class="slide">
          {slide_header("العرض والطلب: بيع وشراء", "مقارنة المعروض للبيع مع طلبات الشراء")}
          <main class="sale-page">
            <div class="split-panels">
              <div class="panel">
                <h2>حجم العرض والطلب</h2>
                {bar_chart_html([
                    {"label": "عروض للبيع", "value": metrics["sale_count"]},
                    {"label": "طلبات شراء", "value": metrics["buy_count"]},
                ], "سجل", color=css_hex(BRAND["blue"]))}
              </div>
              <div class="panel">
                <h2>أعلى مناطق عروض البيع</h2>
                {bar_chart_html(metrics["sale_city_top"], "عرض", color=css_hex(BRAND["green"]))}
              </div>
            </div>
            <div class="kpi-row">
              {kpi_html("نسبة العرض للطلب", fmt_decimal(ratio, 1), "عرض بيع لكل طلب شراء")}
              {kpi_html("متوسط أسعار البيع المعلنة", fmt_kd(sale_stats["avg"]), "د.ك")}
              {kpi_html("طلبات شراء بسعر/ميزانية معلنة", fmt_int(buy_stats["priced_count"]), "طلب")}
            </div>
          </main>
          {source_footer(4)}
        </section>
        """
    )

    sections.append(
        f"""
        <section class="slide">
          {slide_header("البيوت المعروضة للبيع", "توزيع العروض والأسعار المعلنة")}
          <main class="listing-main">
            <div class="two-col wide-left">
              <div class="panel">
                <h2>توزيع البيوت للبيع حسب المنطقة</h2>
                {bar_chart_html(metrics["house_city_top"], "عرض", color=css_hex(BRAND["gold"]))}
              </div>
              <div class="side-kpis">
                {kpi_html("عدد البيوت للبيع", fmt_int(metrics["house_sale_count"]), "عرض")}
                {kpi_html("متوسط السعر المعلن", fmt_kd(house_stats["avg"]), "د.ك")}
                {kpi_html("وسيط السعر المعلن", fmt_kd(house_stats["median"]), "د.ك")}
              </div>
            </div>
            <div class="panel low-panel">
              <h2>شرائح السعر</h2>
              {band_cards_html(metrics["house_price_bands"], "عرض")}
            </div>
          </main>
          {source_footer(5)}
        </section>
        """
    )

    sections.append(
        f"""
        <section class="slide">
          {slide_header("الإيجارات", "عروض الإيجار وطلبات الإيجار")}
          <main class="listing-main">
            <div class="two-col wide-left">
              <div class="panel">
                <h2>أعلى مناطق عروض الإيجار</h2>
                {bar_chart_html(metrics["rent_city_top"], "عرض", color=css_hex(BRAND["teal"]))}
              </div>
              <div class="side-kpis">
                {kpi_html("عروض الإيجار", fmt_int(metrics["rent_count"]), "عرض")}
                {kpi_html("طلبات الإيجار", fmt_int(metrics["rent_request_count"]), "طلب")}
                {kpi_html("متوسط الإيجار المعلن", fmt_kd(rent_stats["avg"]), "د.ك شهرياً")}
                {kpi_html("وسيط الإيجار المعلن", fmt_kd(rent_stats["median"]), "د.ك شهرياً")}
              </div>
            </div>
            <div class="panel low-panel">
              <h2>شرائح الإيجار الشهري</h2>
              {band_cards_html(metrics["rent_price_bands"], "عرض")}
            </div>
          </main>
          {source_footer(6)}
        </section>
        """
    )

    sections.append(
        f"""
        <section class="slide">
          {slide_header("الإحصائيات حسب المحافظات", "توزيع الحركة والبيع والشراء والإيجارات")}
          <main>
            <div class="panel gov-full-panel">
              <h2>ملخص المحافظات</h2>
              {governorate_table_html(metrics["governorate_summary"])}
            </div>
          </main>
          {source_footer(7)}
        </section>
        """
    )

    sections.append(
        f"""
        <section class="slide">
          {slide_header("المصادر والمنهجية", "مصادر الأرقام وحدود الاستخدام")}
          <main>
            <ul class="methodology">
              <li>مصدر الأرقام: واجهة البحث العامة لمنصة الفريج، مع إثراء التصنيف من التفاصيل المتاحة لكل إعلان.</li>
              <li>تاريخ السحب والتحليل: {html.escape(SNAPSHOT_LABEL)}.</li>
              <li>احتساب المتوسطات تم من الأسعار المعلنة فقط؛ والقيم الصفرية عوملت كـ “سعر غير معلن”.</li>
              <li>تم استبعاد بيانات التواصل وبيانات الملاك والعملاء من التقرير، واعتماد المجاميع والمؤشرات فقط.</li>
              <li>تم عرض الإحصائيات حسب المحافظات عند توفر المحافظة في بيانات الإعلان.</li>
            </ul>
          </main>
          {source_footer(8)}
        </section>
        """
    )

    css = f"""
    @page {{ size: 13.333in 7.5in; margin: 0; }}
    * {{ box-sizing: border-box; }}
    body {{ margin: 0; background: #fff; direction: rtl; font-family: Arial, Tahoma, sans-serif; color: {css_hex(BRAND["ink"])}; }}
    .slide {{ position: relative; width: 13.333in; height: 7.5in; overflow: hidden; page-break-after: always; background: {css_hex(BRAND["soft"])}; padding: 0.48in 0.62in; }}
    .slide::before {{ content: ""; position: absolute; inset: 0 0 auto 0; height: 0.17in; background: {css_hex(BRAND["dark"])}; }}
    header {{ height: 0.95in; display: flex; align-items: flex-start; justify-content: space-between; gap: 0.4in; }}
    header h1 {{ margin: 0.02in 0 0; font-size: 31px; line-height: 1.18; letter-spacing: 0; }}
    header p {{ margin: 0.08in 0 0; font-size: 17px; color: {css_hex(BRAND["muted"])}; font-weight: 700; }}
    .logo {{ width: 2.35in; height: auto; object-fit: contain; }}
    main {{ position: relative; height: 5.72in; }}
    footer {{ position: absolute; left: 0.62in; right: 0.62in; bottom: 0.24in; border-top: 1px solid {css_hex(BRAND["line"])}; padding-top: 0.06in; display: flex; justify-content: space-between; color: {css_hex(BRAND["muted"])}; font-size: 10px; font-weight: 700; }}
    .cover {{ background: linear-gradient(90deg, rgba(15, 23, 42, 0.16), rgba(15, 23, 42, 0.72) 45%, rgba(15, 23, 42, 0.96) 100%), url("../assets/kuwait_glass_cover.png") center/cover no-repeat; color: white; padding: 0.55in 0.76in; }}
    .cover::before {{ background: {css_hex(BRAND["gold"])}; }}
    .cover-logo {{ position: absolute; top: 0.48in; right: 0.92in; width: 2.8in; }}
    .cover-glass {{ position: absolute; right: 0.86in; top: 1.28in; width: 6.75in; padding: 0.42in 0.52in 0.38in; background: rgba(15, 23, 42, 0.76); border: 1px solid rgba(226, 232, 240, 0.28); border-right: 7px solid {css_hex(BRAND["gold"])}; border-radius: 8px; box-shadow: 0 28px 70px rgba(0, 0, 0, 0.34); backdrop-filter: blur(10px); }}
    .cover h1 {{ margin: 0; font-size: 39px; line-height: 1.18; letter-spacing: 0; }}
    .cover h2 {{ margin: 0.18in 0 0; font-size: 21px; color: {css_hex(BRAND["gold_2"])}; }}
    .cover p {{ margin: 0.16in 0 0; color: #e5e7eb; font-size: 17px; line-height: 1.55; }}
    .cover-kicker {{ color: #dbeafe !important; font-weight: 700; }}
    .cover-chips {{ display: grid; grid-template-columns: repeat(4, 1fr); gap: 0.12in; margin-top: 0.3in; }}
    .cover-chips span {{ display: block; border: 1px solid rgba(203, 213, 225, 0.28); border-radius: 7px; padding: 0.13in 0.08in; background: rgba(30, 41, 59, 0.82); color: #dbeafe; font-size: 13px; text-align: center; font-weight: 700; }}
    .cover-chips b {{ display: block; color: white; font-size: 23px; margin-bottom: 0.03in; }}
    .kpi-grid {{ display: grid; gap: 0.22in; }}
    .kpi-grid.six {{ grid-template-columns: repeat(3, 1fr); padding-top: 0.3in; }}
    .kpi-row {{ display: grid; grid-template-columns: repeat(3, 1fr); gap: 0.18in; margin-top: 0.22in; }}
    .kpi {{ background: #fff; border: 1px solid {css_hex(BRAND["line"])}; border-right: 6px solid {css_hex(BRAND["gold"])}; border-radius: 7px; min-height: 1.18in; padding: 0.15in 0.18in; }}
    .kpi span {{ display: block; color: {css_hex(BRAND["muted"])}; font-size: 15px; font-weight: 700; }}
    .kpi strong {{ display: block; margin-top: 0.06in; font-size: 34px; line-height: 1.0; color: {css_hex(BRAND["ink"])}; }}
    .kpi em {{ display: block; margin-top: 0.06in; font-size: 13px; color: {css_hex(BRAND["muted"])}; font-style: normal; font-weight: 700; }}
    .kpi small {{ display: block; margin-top: 0.08in; font-size: 12px; color: {css_hex(BRAND["muted"])}; font-weight: 700; line-height: 1.4; }}
    .two-col {{ display: grid; grid-template-columns: 1fr 3.2in; gap: 0.26in; padding-top: 0.23in; align-items: start; }}
    .two-col.wide-left {{ grid-template-columns: 1fr 3.05in; }}
    .split-panels {{ display: grid; grid-template-columns: 1fr 1fr; gap: 0.24in; padding-top: 0.25in; }}
    .compact-panels {{ padding-top: 0.12in; }}
    .compact-panels .panel {{ min-height: 2.38in; }}
    .side-kpis {{ display: grid; gap: 0.14in; }}
    .side-kpis .kpi {{ min-height: 1.05in; }}
    .panel {{ background: #fff; border: 1px solid {css_hex(BRAND["line"])}; border-radius: 7px; padding: 0.2in; min-height: 3.75in; }}
    .panel h2 {{ margin: 0 0 0.15in; font-size: 20px; line-height: 1.25; }}
    .bars {{ display: grid; gap: 0.12in; padding-top: 0.04in; }}
    .bar-row {{ display: grid; grid-template-columns: 1.6in 1fr 1.05in; gap: 0.15in; align-items: center; min-height: 0.34in; }}
    .bar-label {{ font-size: 16px; color: {css_hex(BRAND["ink"])}; font-weight: 700; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }}
    .bar-track {{ height: 0.15in; background: #e2e8f0; border-radius: 2px; position: relative; overflow: hidden; }}
    .bar-track i {{ position: absolute; top: 0; right: 0; height: 100%; border-radius: 2px; }}
    .bar-row strong {{ font-size: 15px; color: {css_hex(BRAND["muted"])}; text-align: left; white-space: nowrap; }}
    .low-panel {{ margin-top: 0.16in; min-height: 1.12in; padding: 0.13in 0.18in; }}
    .band-grid {{ display: grid; grid-template-columns: repeat(6, 1fr); gap: 0.1in; }}
    .band-card {{ min-height: 0.66in; border: 1px solid #dbe3ef; border-right: 4px solid {css_hex(BRAND["blue"])}; border-radius: 6px; padding: 0.08in 0.08in; background: #f8fafc; }}
    .band-card span {{ display: block; font-size: 12px; line-height: 1.25; color: {css_hex(BRAND["muted"])}; font-weight: 700; white-space: normal; }}
    .band-card strong {{ display: inline-block; margin-top: 0.03in; font-size: 20px; line-height: 1; color: {css_hex(BRAND["ink"])}; }}
    .band-card em {{ display: inline-block; margin-right: 0.04in; font-size: 10px; color: {css_hex(BRAND["muted"])}; font-style: normal; font-weight: 700; }}
    .table-panel {{ margin-top: 0.16in; min-height: 2.5in; padding: 0.14in 0.18in; }}
    .gov-full-panel {{ margin-top: 0.34in; min-height: 4.55in; padding: 0.22in 0.24in; }}
    .gov-full-panel h2 {{ font-size: 23px; margin-bottom: 0.24in; }}
    .gov-table {{ width: 100%; border-collapse: collapse; table-layout: fixed; font-size: 15px; }}
    .gov-table th {{ background: {css_hex(BRAND["dark"])}; color: #fff; padding: 0.1in 0.05in; border: 1px solid {css_hex(BRAND["dark"])}; font-weight: 700; }}
    .gov-table td {{ background: #fff; padding: 0.085in 0.05in; border: 1px solid {css_hex(BRAND["line"])}; text-align: center; font-weight: 700; color: {css_hex(BRAND["ink"])}; }}
    .gov-table td:first-child {{ text-align: right; color: {css_hex(BRAND["muted"])}; }}
    .priority p {{ margin: 0 0 0.16in; font-size: 19px; line-height: 1.58; color: {css_hex(BRAND["ink"])}; }}
    .methodology {{ margin: 0.32in 0.18in 0 0; padding: 0; list-style-position: inside; font-size: 23px; line-height: 1.78; font-weight: 700; }}
    .source-box {{ margin-top: 0.34in; background: #fff; border: 1px solid {css_hex(BRAND["line"])}; border-right: 6px solid {css_hex(BRAND["gold"])}; border-radius: 7px; padding: 0.2in 0.25in; }}
    .source-box strong {{ font-size: 20px; color: {css_hex(BRAND["ink"])}; }}
    .source-box p {{ margin: 0.08in 0 0; font-size: 18px; line-height: 1.6; color: {css_hex(BRAND["muted"])}; font-weight: 700; }}
    .sale-page .split-panels {{ padding-top: 0.18in; }}
    .sale-page .panel {{ min-height: 3.18in; }}
    .sale-page .kpi-row {{ margin-top: 0.16in; }}
    .sale-page .kpi-row .kpi {{ min-height: 0.9in; padding: 0.1in 0.16in; }}
    .sale-page .kpi-row .kpi strong {{ font-size: 30px; }}
    .listing-main .two-col {{ padding-top: 0.18in; }}
    .listing-main .two-col > .panel {{ min-height: 3.28in; }}
    .listing-main .side-kpis {{ gap: 0.08in; }}
    .listing-main .side-kpis .kpi {{ min-height: 0.76in; padding: 0.09in 0.16in; }}
    .listing-main .side-kpis .kpi span {{ font-size: 14px; }}
    .listing-main .side-kpis .kpi strong {{ font-size: 28px; }}
    .listing-main .side-kpis .kpi em {{ font-size: 12px; margin-top: 0.04in; }}
    .listing-main .low-panel {{ margin-top: 0.14in; }}
    .insight {{ margin: 0.18in 0 0; padding: 0.12in 0.18in; background: #fff; border-right: 6px solid {css_hex(BRAND["gold"])}; color: {css_hex(BRAND["ink"])}; font-size: 17px; line-height: 1.38; font-weight: 700; }}
    .empty {{ font-size: 16px; color: {css_hex(BRAND["muted"])}; padding: 0.1in; }}
    """
    html_doc = f"""<!doctype html>
<html lang="ar" dir="rtl">
<head>
  <meta charset="utf-8">
  <title>{html.escape(REPORT_TITLE)}</title>
  <style>{css}</style>
</head>
<body>
  {"".join(sections)}
</body>
</html>
"""
    REPORT_HTML_PATH.write_text(html_doc, encoding="utf-8")


def apply_arabic_text_direction(paragraph) -> None:
    paragraph._p.get_or_add_pPr().set("rtl", "1")


def apply_arabic_run_language(run) -> None:
    r_pr = run._r.get_or_add_rPr()
    r_pr.set("lang", "ar-KW")
    r_pr.set("altLang", "en-US")


def add_rect(slide, left: float, top: float, width: float, height: float, fill: str, line: str | None = None) -> object:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill)
    shape.line.color.rgb = hex_to_rgb(line or fill)
    shape.line.width = Pt(0.5)
    return shape


def add_round_rect(slide, left: float, top: float, width: float, height: float, fill: str, line: str | None = None) -> object:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill)
    shape.line.color.rgb = hex_to_rgb(line or fill)
    shape.line.width = Pt(0.7)
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: int = 20,
    bold: bool = False,
    color: str = BRAND["ink"],
    align: PP_ALIGN = PP_ALIGN.RIGHT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> object:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    apply_arabic_text_direction(paragraph)
    run = paragraph.add_run()
    run.text = text
    apply_arabic_run_language(run)
    run.font.name = FONT_AR
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)
    return box


def add_multiline(
    slide,
    lines: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: int = 18,
    color: str = BRAND["ink"],
    bold: bool = False,
    spacing_after: int = 7,
) -> object:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph.space_after = Pt(spacing_after)
        apply_arabic_text_direction(paragraph)
        run = paragraph.add_run()
        run.text = line
        apply_arabic_run_language(run)
        run.font.name = FONT_AR
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(color)
    return box


def add_logo(slide, left: float, top: float, width: float) -> None:
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(left), Inches(top), width=Inches(width))


def add_slide_header(slide, title: str, subtitle: str, page: int) -> None:
    add_rect(slide, 0, 0, 13.333, 0.34, BRAND["dark"])
    add_text(slide, title, 6.0, 0.52, 6.5, 0.48, font_size=26, bold=True)
    add_text(slide, subtitle, 6.0, 1.0, 6.5, 0.34, font_size=14, bold=True, color=BRAND["muted"])
    add_logo(slide, 0.55, 0.48, 2.35)
    add_rect(slide, 12.68, 0.52, 0.12, 0.72, BRAND["gold"])
    add_footer(slide, page)


def add_footer(slide, page: int) -> None:
    add_rect(slide, 0.55, 7.02, 12.25, 0.01, BRAND["line"])
    add_text(slide, f"{CLIENT_NAME} | المصدر: {SOURCE_NAME}", 4.2, 7.08, 8.6, 0.25, font_size=8, bold=True, color=BRAND["muted"])
    add_text(slide, f"{page:02d}", 0.55, 7.08, 0.42, 0.25, font_size=8, bold=True, color=BRAND["muted"], align=PP_ALIGN.LEFT)


def add_kpi_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    value: str,
    unit: str,
    note: str = "",
    accent: str = BRAND["gold"],
) -> None:
    add_round_rect(slide, left, top, width, height, BRAND["white"], BRAND["line"])
    compact = height < 1.2
    add_rect(slide, left + width - 0.08, top + (0.12 if compact else 0.16), 0.045, height - (0.24 if compact else 0.32), accent)
    title_size = 10 if compact else 12
    value_size = 22 if compact else 29
    detail_size = 8 if compact else 10
    title_top = 0.1 if compact else 0.13
    value_top = 0.36 if compact else 0.5
    value_height = 0.38 if compact else 0.42
    detail_top = height - (0.24 if compact else 0.32)
    add_text(slide, title, left + 0.18, top + title_top, width - 0.38, 0.22 if compact else 0.28, font_size=title_size, bold=True, color=BRAND["muted"])
    add_text(slide, value, left + 0.18, top + value_top, width - 0.38, value_height, font_size=value_size, bold=True)
    detail = f"{unit} | {note}" if note else unit
    if detail:
        add_text(slide, detail, left + 0.18, top + detail_top, width - 0.38, 0.18 if compact else 0.24, font_size=detail_size, bold=True, color=BRAND["muted"])


def add_bar_chart(slide, data: list[dict], left: float, top: float, width: float, height: float, *, suffix: str = "", color: str = BRAND["blue"]) -> None:
    if not data:
        add_text(slide, "لا توجد قيم قابلة للعرض", left, top, width, 0.32, font_size=14, bold=True, color=BRAND["muted"])
        return
    max_value = max(item["value"] for item in data) or 1
    row_h = height / max(len(data), 1)
    for idx, item in enumerate(data):
        y = top + idx * row_h
        label = str(item["label"])
        value = float(item["value"])
        add_text(slide, label, left + width - 1.78, y, 1.78, 0.28, font_size=11, bold=True)
        track_left = left + 0.95
        track_w = width - 2.9
        add_rect(slide, track_left, y + 0.1, track_w, 0.12, "E2E8F0")
        fill_w = max(0.035, track_w * value / max_value)
        add_rect(slide, track_left + track_w - fill_w, y + 0.1, fill_w, 0.12, color)
        add_text(slide, f"{fmt_int(value)} {suffix}".strip(), left, y, 0.85, 0.28, font_size=10, bold=True, color=BRAND["muted"], align=PP_ALIGN.LEFT)


def add_band_cards(slide, data: list[dict], left: float, top: float, width: float, height: float, *, suffix: str = "", accent: str = BRAND["blue"]) -> None:
    if not data:
        add_text(slide, "لا توجد قيم قابلة للعرض", left, top, width, 0.32, font_size=14, bold=True, color=BRAND["muted"])
        return
    cols = min(6, len(data))
    gap = 0.12
    card_w = (width - gap * (cols - 1)) / cols
    card_h = height
    for idx, item in enumerate(data[:cols]):
        x = left + idx * (card_w + gap)
        add_round_rect(slide, x, top, card_w, card_h, "F8FAFC", BRAND["line"])
        add_rect(slide, x + card_w - 0.04, top + 0.08, 0.03, card_h - 0.16, accent)
        add_text(slide, str(item["label"]), x + 0.08, top + 0.07, card_w - 0.18, 0.3, font_size=10, bold=True, color=BRAND["muted"])
        add_text(slide, fmt_int(item["value"]), x + 0.08, top + 0.42, card_w - 0.18, 0.34, font_size=18, bold=True)


def add_panel(slide, left: float, top: float, width: float, height: float, title: str) -> None:
    add_round_rect(slide, left, top, width, height, BRAND["white"], BRAND["line"])
    add_text(slide, title, left + 0.18, top + 0.16, width - 0.36, 0.34, font_size=17, bold=True)


def add_governorate_table(slide, rows: list[dict], left: float, top: float, width: float, height: float) -> None:
    headers = ["المحافظة", "إجمالي", "للبيع", "شراء", "بيوت", "إيجار", "طلب إيجار"]
    data = [
        [
            row["label"],
            fmt_int(row["total"]),
            fmt_int(row["sale"]),
            fmt_int(row["buy"]),
            fmt_int(row["house_sale"]),
            fmt_int(row["rent"]),
            fmt_int(row["rent_request"]),
        ]
        for row in rows
    ]
    table_shape = slide.shapes.add_table(len(data) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    col_widths = [2.05, 1.15, 1.0, 1.0, 1.05, 1.0, 1.15]
    scale = width / sum(col_widths)
    for idx, col_width in enumerate(col_widths):
        table.columns[idx].width = Inches(col_width * scale)
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(BRAND["dark"])
        cell.margin_left = Inches(0.03)
        cell.margin_right = Inches(0.03)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        apply_arabic_text_direction(paragraph)
        run = paragraph.add_run()
        run.text = header
        apply_arabic_run_language(run)
        run.font.name = FONT_AR
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = hex_to_rgb(BRAND["white"])
    for row_idx, row in enumerate(data, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(BRAND["white"])
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.RIGHT if col_idx == 0 else PP_ALIGN.CENTER
            apply_arabic_text_direction(paragraph)
            run = paragraph.add_run()
            run.text = str(value)
            apply_arabic_run_language(run)
            run.font.name = FONT_AR
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = hex_to_rgb(BRAND["ink"] if col_idx else BRAND["muted"])


def create_pptx(metrics: dict) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    if COVER_BG_PATH.exists():
        slide.shapes.add_picture(str(COVER_BG_PATH), Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
    else:
        add_rect(slide, 0, 0, 13.333, 7.5, BRAND["dark"])
    add_rect(slide, 0, 0, 13.333, 0.12, BRAND["gold"])
    add_round_rect(slide, 6.05, 1.25, 6.55, 4.95, "111827", "334155")
    add_rect(slide, 12.48, 1.55, 0.07, 4.38, BRAND["gold"])
    add_logo(slide, 9.4, 0.48, 2.85)
    add_text(slide, REPORT_TITLE, 6.55, 1.92, 5.55, 0.7, font_size=34, bold=True, color=BRAND["white"])
    add_text(slide, REPORT_SUBTITLE, 6.55, 2.78, 5.55, 0.42, font_size=18, bold=True, color=BRAND["gold_2"])
    add_text(slide, SNAPSHOT_LABEL, 6.55, 3.34, 5.55, 0.34, font_size=15, bold=True, color="E5E7EB")
    add_kpi_card(slide, 6.55, 4.18, 1.3, 1.05, "الحركة", fmt_int(metrics["total_core"]), "", "", BRAND["gold"])
    add_kpi_card(slide, 7.95, 4.18, 1.3, 1.05, "للبيع", fmt_int(metrics["sale_count"]), "", "", BRAND["blue"])
    add_kpi_card(slide, 9.35, 4.18, 1.3, 1.05, "إيجار", fmt_int(metrics["rent_count"]), "", "", BRAND["teal"])
    add_kpi_card(slide, 10.75, 4.18, 1.3, 1.05, "شراء", fmt_int(metrics["buy_count"]), "", "", BRAND["green"])
    add_footer(slide, 1)

    # 2. Executive summary
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "الملخص التنفيذي", "الأرقام التي تجيب مباشرة عن المطلوب", 2)
    cards = [
        ("حركة الدلال", fmt_int(metrics["total_core"]), "إعلان/طلب منشور", f"{fmt_int(metrics['details_loaded'])} إعلان بتفاصيل متاحة", BRAND["gold"]),
        ("عروض البيع", fmt_int(metrics["sale_count"]), "عرض", f"{fmt_int(metrics['sale_stats']['priced_count'])} بسعر معلن", BRAND["blue"]),
        ("طلبات الشراء", fmt_int(metrics["buy_count"]), "طلب", f"{fmt_decimal(metrics['sale_to_buy_ratio'])} عرض لكل طلب", BRAND["green"]),
        ("بيوت للبيع", fmt_int(metrics["house_sale_count"]), "عرض", f"متوسط {fmt_kd(metrics['house_stats']['avg'])} د.ك", BRAND["gold"]),
        ("عروض الإيجار", fmt_int(metrics["rent_count"]), "عرض", f"متوسط {fmt_kd(metrics['rent_stats']['avg'])} د.ك", BRAND["teal"]),
        ("طلبات الإيجار", fmt_int(metrics["rent_request_count"]), "طلب", f"نسبة {fmt_decimal(metrics['rent_offer_to_request_ratio'])} مرة", BRAND["red"]),
    ]
    x_positions = [0.75, 4.55, 8.35]
    y_positions = [1.65, 3.3]
    for idx, (title, value, unit, note, accent) in enumerate(cards):
        add_kpi_card(slide, x_positions[idx % 3], y_positions[idx // 3], 3.2, 1.35, title, value, unit, note, accent)

    # 3. Broker movement
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "حركة الدلال", "قياس النشاط المنشور على المنصة", 3)
    add_panel(slide, 0.75, 1.55, 5.75, 3.55, "الحركة حسب نوع المعاملة")
    add_bar_chart(slide, metrics["type_counts"], 1.0, 2.18, 5.25, 1.95, suffix="سجل", color=BRAND["gold"])
    add_panel(slide, 6.85, 1.55, 5.3, 3.55, "أعلى مناطق الحركة")
    add_bar_chart(slide, metrics["movement_city_top"][:6], 7.1, 2.18, 4.8, 1.95, suffix="سجل", color=BRAND["blue"])
    add_kpi_card(slide, 0.75, 5.15, 3.55, 1.08, "إجمالي الحركة", fmt_int(metrics["total_core"]), "", "", BRAND["gold"])
    add_kpi_card(slide, 4.7, 5.15, 3.55, 1.08, "تفاصيل متاحة", fmt_int(metrics["details_loaded"]), "", "", BRAND["blue"])
    add_kpi_card(slide, 8.65, 5.15, 3.5, 1.08, "إعلانات مباشرة", fmt_int(metrics["direct_count"]), "", "", BRAND["green"])

    # 4. Sale and purchase
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "العرض والطلب: بيع وشراء", "مقارنة المعروض للبيع مع طلبات الشراء", 4)
    add_panel(slide, 0.75, 1.55, 5.6, 3.2, "حجم العرض والطلب")
    add_bar_chart(
        slide,
        [
            {"label": "عروض للبيع", "value": metrics["sale_count"]},
            {"label": "طلبات شراء", "value": metrics["buy_count"]},
        ],
        1.0,
        2.25,
        5.05,
        1.0,
        suffix="سجل",
        color=BRAND["blue"],
    )
    add_panel(slide, 6.65, 1.55, 5.5, 3.2, "أعلى مناطق عروض البيع")
    add_bar_chart(slide, metrics["sale_city_top"][:5], 6.9, 2.12, 5.0, 1.95, suffix="عرض", color=BRAND["green"])
    add_kpi_card(slide, 0.75, 5.05, 3.55, 1.1, "نسبة العرض للطلب", fmt_decimal(metrics["sale_to_buy_ratio"], 1), "", "", BRAND["gold"])
    add_kpi_card(slide, 4.7, 5.05, 3.55, 1.1, "متوسط البيع المعلن", fmt_kd(metrics["sale_stats"]["avg"]), "", "", BRAND["blue"])
    add_kpi_card(slide, 8.65, 5.05, 3.5, 1.1, "طلبات شراء بسعر معلن", fmt_int(metrics["buy_stats"]["priced_count"]), "", "", BRAND["green"])

    # 5. House sale
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "البيوت المعروضة للبيع", "توزيع العروض والأسعار المعلنة", 5)
    add_panel(slide, 0.75, 1.55, 7.15, 3.25, "توزيع البيوت للبيع حسب المنطقة")
    add_bar_chart(slide, metrics["house_city_top"][:5], 1.0, 2.13, 6.65, 2.0, suffix="عرض", color=BRAND["gold"])
    add_kpi_card(slide, 8.3, 1.55, 3.85, 0.94, "عدد البيوت للبيع", fmt_int(metrics["house_sale_count"]), "", "", BRAND["gold"])
    add_kpi_card(slide, 8.3, 2.62, 3.85, 0.94, "متوسط السعر المعلن", fmt_kd(metrics["house_stats"]["avg"]), "", "", BRAND["blue"])
    add_kpi_card(slide, 8.3, 3.69, 3.85, 0.94, "وسيط السعر المعلن", fmt_kd(metrics["house_stats"]["median"]), "", "", BRAND["green"])
    add_panel(slide, 0.75, 4.95, 11.4, 1.45, "شرائح السعر")
    add_band_cards(slide, metrics["house_price_bands"], 1.0, 5.44, 10.9, 0.88, suffix="عرض", accent=BRAND["blue"])

    # 6. Rent
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "الإيجارات", "عروض الإيجار وطلبات الإيجار", 6)
    add_panel(slide, 0.75, 1.55, 7.15, 3.25, "أعلى مناطق عروض الإيجار")
    add_bar_chart(slide, metrics["rent_city_top"][:5], 1.0, 2.13, 6.65, 2.0, suffix="عرض", color=BRAND["teal"])
    add_kpi_card(slide, 8.3, 1.55, 3.85, 0.94, "عروض الإيجار", fmt_int(metrics["rent_count"]), "", "", BRAND["teal"])
    add_kpi_card(slide, 8.3, 2.62, 3.85, 0.94, "طلبات الإيجار", fmt_int(metrics["rent_request_count"]), "", "", BRAND["red"])
    add_kpi_card(slide, 8.3, 3.69, 3.85, 0.94, "متوسط الإيجار المعلن", fmt_kd(metrics["rent_stats"]["avg"]), "", "", BRAND["green"])
    add_panel(slide, 0.75, 4.95, 11.4, 1.45, "شرائح الإيجار الشهري")
    add_band_cards(slide, metrics["rent_price_bands"], 1.0, 5.44, 10.9, 0.88, suffix="عرض", accent=BRAND["green"])

    # 7. Governorates
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "الإحصائيات حسب المحافظات", "توزيع الحركة والبيع والشراء والإيجارات", 7)
    add_panel(slide, 0.75, 1.55, 11.4, 4.95, "ملخص المحافظات")
    add_governorate_table(slide, metrics["governorate_summary"], 1.0, 2.12, 10.9, 3.75)

    # 8. Sources
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "المصادر والمنهجية", "مصادر الأرقام وحدود الاستخدام", 8)
    lines = [
        f"مصدر الأرقام: {SOURCE_NAME}، مع إثراء التصنيف من التفاصيل المتاحة لكل إعلان.",
        f"تاريخ السحب والتحليل: {SNAPSHOT_LABEL}.",
        "احتساب المتوسطات تم من الأسعار المعلنة فقط؛ والقيم الصفرية عوملت كسعر غير معلن.",
        "تم استبعاد بيانات التواصل وبيانات الملاك والعملاء من التقرير، واعتماد المجاميع والمؤشرات فقط.",
        "تم عرض الإحصائيات حسب المحافظات عند توفر المحافظة في بيانات الإعلان.",
    ]
    add_multiline(slide, lines, 1.0, 1.78, 11.2, 3.75, font_size=20, bold=True, spacing_after=12)

    prs.save(PPTX_PATH)


def style_header(ws, row: int = 1) -> None:
    fill = PatternFill("solid", fgColor=BRAND["dark"])
    font = Font(name="Arial", size=11, bold=True, color=BRAND["white"])
    border = Border(
        left=Side(style="thin", color=BRAND["line"]),
        right=Side(style="thin", color=BRAND["line"]),
        top=Side(style="thin", color=BRAND["line"]),
        bottom=Side(style="thin", color=BRAND["line"]),
    )
    for cell in ws[row]:
        cell.fill = fill
        cell.font = font
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True, readingOrder=2)
        cell.border = border
    ws.row_dimensions[row].height = 30


def style_body(ws) -> None:
    border = Border(
        left=Side(style="thin", color="E2E8F0"),
        right=Side(style="thin", color="E2E8F0"),
        top=Side(style="thin", color="E2E8F0"),
        bottom=Side(style="thin", color="E2E8F0"),
    )
    for row in ws.iter_rows(min_row=2, max_row=ws.max_row, max_col=ws.max_column):
        for cell in row:
            cell.alignment = Alignment(horizontal="right", vertical="top", wrap_text=True, readingOrder=2)
            cell.border = border
            cell.font = Font(name="Arial", size=10, color=BRAND["ink"])


def set_widths(ws, widths: list[int]) -> None:
    for idx, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(idx)].width = width
    ws.freeze_panes = "A2"
    ws.sheet_view.rightToLeft = True
    ws.auto_filter.ref = ws.dimensions


def write_sheet(ws, headers: list[str], rows: list[list[object]], widths: list[int]) -> None:
    ws.append(headers)
    for row in rows:
        ws.append(row)
    style_header(ws)
    style_body(ws)
    set_widths(ws, widths)
    ws.sheet_properties.pageSetUpPr.fitToPage = True
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 0


def listing_rows(rows: list[dict]) -> list[list[object]]:
    result = []
    for row in rows:
        result.append(
            [
                row["code"],
                row["transaction_type_name"],
                row["property_type_name"],
                row["detail_class"],
                row["governorate_name"],
                row["city_name"],
                row["price"] if row["is_priced"] else "",
                "نعم" if row["is_priced"] else "غير معلن",
                row["area"] or "",
                row["listing_mode"],
                row["condition"],
                row["features"],
                "نعم" if row["has_image"] else "لا",
                row["published_date"].isoformat() if row.get("published_date") else "",
                row["days_since"] if row.get("days_since") is not None else "",
                "نعم" if row["detail_loaded"] else "لا",
                row["source_url"],
            ]
        )
    return result


def create_excel(metrics: dict) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "ملخص"

    summary_rows = [
        ["حركة الدلال", "إجمالي الحركة ضمن نطاق الطلب", metrics["total_core"], "سجل", SOURCE_NAME, "يشمل البيع، الإيجار، طلب الشراء، طلب الإيجار"],
        ["حركة الدلال", "إعلانات بتفاصيل متاحة", metrics["details_loaded"], "إعلان/طلب", SOURCE_NAME, "تم إثراء السجل من تفاصيل الإعلان عندما كانت متاحة"],
        ["حركة الدلال", "إعلانات مباشرة", metrics["direct_count"], "إعلان/طلب", SOURCE_NAME, "بحسب عبارة نوع الإعلان في تفاصيل الإعلان عند توفرها"],
        ["البيع والشراء", "عروض للبيع", metrics["sale_count"], "عرض", SOURCE_NAME, ""],
        ["البيع والشراء", "طلبات شراء", metrics["buy_count"], "طلب", SOURCE_NAME, ""],
        ["البيع والشراء", "نسبة عروض البيع إلى طلبات الشراء", metrics["sale_to_buy_ratio"], "مرة", SOURCE_NAME, ""],
        ["البيع", "متوسط أسعار البيع المعلنة", metrics["sale_stats"]["avg"] or "", "د.ك", SOURCE_NAME, "من الأسعار المعلنة فقط"],
        ["البيوت للبيع", "عدد البيوت المعروضة للبيع", metrics["house_sale_count"], "عرض", SOURCE_NAME, "بيت/فيلا/قسيمة/منزل"],
        ["البيوت للبيع", "متوسط سعر البيوت المعلن", metrics["house_stats"]["avg"] or "", "د.ك", SOURCE_NAME, "من الأسعار المعلنة فقط"],
        ["البيوت للبيع", "وسيط سعر البيوت المعلن", metrics["house_stats"]["median"] or "", "د.ك", SOURCE_NAME, "من الأسعار المعلنة فقط"],
        ["الإيجارات", "عروض الإيجار", metrics["rent_count"], "عرض", SOURCE_NAME, ""],
        ["الإيجارات", "طلبات الإيجار", metrics["rent_request_count"], "طلب", SOURCE_NAME, ""],
        ["الإيجارات", "متوسط الإيجار الشهري المعلن", metrics["rent_stats"]["avg"] or "", "د.ك/شهر", SOURCE_NAME, "من الأسعار المعلنة فقط"],
    ]
    write_sheet(
        ws,
        ["المحور", "المؤشر", "القيمة", "الوحدة", "المصدر", "ملاحظات"],
        summary_rows,
        [18, 38, 16, 16, 32, 44],
    )

    ws_types = wb.create_sheet("حسب_نوع_المعاملة")
    type_rows = []
    for item in metrics["type_counts_all"]:
        type_id = next((tid for tid, name in metrics["type_names"].items() if name == item["label"]), "")
        rows = metrics["by_type"].get(type_id, []) if type_id else []
        stats = price_stats(rows)
        type_rows.append(
            [
                item["label"],
                item["value"],
                stats["priced_count"],
                stats["undisclosed_count"],
                stats["avg"] or "",
                stats["median"] or "",
                pct(item["value"], metrics["total_all"]),
            ]
        )
    write_sheet(
        ws_types,
        ["نوع المعاملة", "عدد السجلات", "أسعار معلنة", "أسعار غير معلنة", "متوسط السعر", "وسيط السعر", "الحصة من الإجمالي %"],
        type_rows,
        [22, 16, 16, 18, 18, 18, 20],
    )

    ws_agents = wb.create_sheet("حركة_الدلال")
    movement_rows = []
    area_labels = sorted({row["city_name"] for row in metrics["core_rows"]})
    for city in area_labels:
        city_rows = [row for row in metrics["core_rows"] if row["city_name"] == city]
        movement_rows.append(
            [
                city,
                len(city_rows),
                sum(1 for row in city_rows if row["transaction_type_id"] == 1),
                sum(1 for row in city_rows if row["transaction_type_id"] == 2),
                sum(1 for row in city_rows if row["transaction_type_id"] == 3),
                sum(1 for row in city_rows if row["transaction_type_id"] == 4),
                sum(1 for row in city_rows if "مباشر" in row.get("listing_mode", "")),
            ]
        )
    movement_rows.sort(key=lambda row: row[1], reverse=True)
    write_sheet(
        ws_agents,
        ["المنطقة", "إجمالي الحركة", "للبيع", "للإيجار", "مطلوب شراء", "مطلوب إيجار", "مباشر"],
        movement_rows,
        [24, 16, 14, 14, 16, 16, 14],
    )

    ws_governorates = wb.create_sheet("حسب_المحافظات")
    governorate_rows = [
        [
            row["label"],
            row["total"],
            row["sale"],
            row["buy"],
            row["house_sale"],
            row["rent"],
            row["rent_request"],
        ]
        for row in metrics["governorate_summary"]
    ]
    write_sheet(
        ws_governorates,
        ["المحافظة", "إجمالي الحركة", "للبيع", "مطلوب شراء", "بيوت للبيع", "إيجار", "طلب إيجار"],
        governorate_rows,
        [22, 16, 12, 16, 15, 12, 14],
    )

    ws_sale_buy = wb.create_sheet("البيع_والشراء")
    area_labels = sorted({row["city_name"] for row in metrics["sale_rows"] + metrics["buy_rows"]})
    sale_buy_rows = []
    for city in area_labels:
        city_sales = [row for row in metrics["sale_rows"] if row["city_name"] == city]
        city_buy = [row for row in metrics["buy_rows"] if row["city_name"] == city]
        if not city_sales and not city_buy:
            continue
        sale_buy_rows.append([city, len(city_sales), len(city_buy), round(len(city_sales) / len(city_buy), 1) if city_buy else "", price_stats(city_sales)["avg"] or ""])
    sale_buy_rows.sort(key=lambda row: (row[1] + row[2]), reverse=True)
    write_sheet(
        ws_sale_buy,
        ["المنطقة", "عروض للبيع", "طلبات شراء", "نسبة العرض للطلب", "متوسط سعر البيع المعلن"],
        sale_buy_rows,
        [24, 16, 16, 18, 22],
    )

    ws_houses = wb.create_sheet("البيوت_للبيع")
    write_sheet(
        ws_houses,
        ["الكود", "نوع المعاملة", "نوع العقار", "التصنيف التفصيلي", "المحافظة", "المنطقة", "السعر", "حالة السعر", "المساحة", "نوع الإعلان", "وصف/حالة مختصرة", "ملامح مستخرجة", "صورة", "تاريخ النشر", "أيام منذ النشر", "تفاصيل متاحة", "صفحة الإعلان"],
        listing_rows(metrics["house_sale_rows"]),
        [14, 16, 18, 20, 22, 20, 16, 16, 12, 16, 28, 28, 10, 16, 16, 16, 60],
    )

    ws_rents = wb.create_sheet("الإيجارات")
    write_sheet(
        ws_rents,
        ["الكود", "نوع المعاملة", "نوع العقار", "التصنيف التفصيلي", "المحافظة", "المنطقة", "السعر", "حالة السعر", "المساحة", "نوع الإعلان", "وصف/حالة مختصرة", "ملامح مستخرجة", "صورة", "تاريخ النشر", "أيام منذ النشر", "تفاصيل متاحة", "صفحة الإعلان"],
        listing_rows(metrics["rent_rows"] + metrics["rent_request_rows"]),
        [14, 16, 18, 20, 22, 20, 16, 16, 12, 16, 28, 28, 10, 16, 16, 16, 60],
    )

    ws_clean = wb.create_sheet("سجل_منظف")
    write_sheet(
        ws_clean,
        ["الكود", "نوع المعاملة", "نوع العقار", "التصنيف التفصيلي", "المحافظة", "المنطقة", "السعر", "حالة السعر", "المساحة", "نوع الإعلان", "وصف/حالة مختصرة", "ملامح مستخرجة", "صورة", "تاريخ النشر", "أيام منذ النشر", "تفاصيل متاحة", "صفحة الإعلان"],
        listing_rows(metrics["core_rows"]),
        [14, 16, 18, 20, 22, 20, 16, 16, 12, 16, 28, 28, 10, 16, 16, 16, 60],
    )

    ws_sources = wb.create_sheet("المصادر_والمنهجية")
    sources_rows = [
        ["الموقع الرسمي للشركة", ALFORAIJ_FRONT_URL, "مرجع الهوية والموقع الرسمي"],
        ["واجهة البحث العامة", ALFORAIJ_SEARCH_URL, "مرجع السجلات المنشورة"],
        ["أنواع المعاملات", TRANSACTION_TYPES_URL, "تصنيف للبيع/للإيجار/مطلوب للشراء/مطلوب للإيجار/للبدل"],
        ["API عروض البيع", f"{LISTINGS_URL}?page=1&pageSize=100&transactionType=1", "مصدر عروض البيع"],
        ["API عروض الإيجار", f"{LISTINGS_URL}?page=1&pageSize=100&transactionType=2", "مصدر عروض الإيجار"],
        ["API طلبات الشراء", f"{LISTINGS_URL}?page=1&pageSize=100&transactionType=3", "مصدر طلبات الشراء"],
        ["API طلبات الإيجار", f"{LISTINGS_URL}?page=1&pageSize=100&transactionType=4", "مصدر طلبات الإيجار"],
        ["منهجية السعر", "", "تم احتساب المتوسط والوسيط من السعر المعلن فقط؛ الصفر أو الفراغ يعامل كسعر غير معلن."],
        ["تفاصيل الإعلان", f"{DETAIL_URL}/{{id}}", "تمت محاولة فحص تفاصيل كل إعلان للتصنيف واستخراج المساحة والملامح، واستخدمت التفاصيل عند توفرها دون عرض بيانات شخصية."],
        ["حماية البيانات", "", "لا يحتوي التقرير أو السجل المنظف على هواتف أو أسماء ملاك/عملاء أو أسماء موظفين."],
    ]
    write_sheet(ws_sources, ["المصدر", "الرابط", "الاستخدام"], sources_rows, [28, 82, 58])

    for wsx in wb.worksheets:
        wsx.sheet_format.defaultRowHeight = 22
    wb.save(EXCEL_PATH)


def create_sources_html(metrics: dict) -> None:
    links = [
        ("الموقع الرسمي للشركة", ALFORAIJ_FRONT_URL),
        ("واجهة البحث العامة", ALFORAIJ_SEARCH_URL),
        ("سجل أنواع المعاملات", TRANSACTION_TYPES_URL),
        ("تفاصيل الإعلان", f"{DETAIL_URL}/{{id}}"),
        ("عروض البيع", f"{LISTINGS_URL}?page=1&pageSize=100&transactionType=1"),
        ("عروض الإيجار", f"{LISTINGS_URL}?page=1&pageSize=100&transactionType=2"),
        ("طلبات الشراء", f"{LISTINGS_URL}?page=1&pageSize=100&transactionType=3"),
        ("طلبات الإيجار", f"{LISTINGS_URL}?page=1&pageSize=100&transactionType=4"),
    ]
    link_html = "\n".join(
        f"""
        <tr>
          <td>{html.escape(title)}</td>
          <td><span dir="ltr">{html.escape(url)}</span></td>
        </tr>
        """
        for title, url in links
    )
    html_doc = f"""<!doctype html>
<html lang="ar" dir="rtl">
<head>
  <meta charset="utf-8">
  <title>مصادر ومنهجية البيانات</title>
  <style>
    @page {{ size: A4; margin: 14mm; }}
    * {{ box-sizing: border-box; }}
    body {{ margin: 0; font-family: Arial, Tahoma, sans-serif; color: {css_hex(BRAND["ink"])}; direction: rtl; }}
    header {{ border-bottom: 5px solid {css_hex(BRAND["gold"])}; padding-bottom: 16px; display: flex; justify-content: space-between; gap: 24px; align-items: flex-start; }}
    header img {{ width: 185px; object-fit: contain; }}
    h1 {{ margin: 0; font-size: 30px; line-height: 1.3; }}
    .sub {{ margin: 8px 0 0; color: {css_hex(BRAND["muted"])}; font-size: 15px; font-weight: 700; }}
    .kpis {{ display: grid; grid-template-columns: repeat(4, 1fr); gap: 10px; margin: 18px 0; }}
    .box {{ border: 1px solid #dbe3ef; border-radius: 8px; padding: 12px; border-right: 5px solid {css_hex(BRAND["gold"])}; }}
    .box span {{ display: block; color: {css_hex(BRAND["muted"])}; font-size: 12px; font-weight: 700; }}
    .box strong {{ display: block; margin-top: 4px; font-size: 24px; }}
    h2 {{ font-size: 20px; margin: 20px 0 10px; }}
    p, li {{ font-size: 15px; line-height: 1.75; font-weight: 700; }}
    table {{ width: 100%; border-collapse: collapse; margin-top: 10px; font-size: 12px; }}
    th, td {{ border: 1px solid #dbe3ef; padding: 9px; vertical-align: top; }}
    th {{ background: {css_hex(BRAND["dark"])}; color: white; }}
    td span[dir="ltr"] {{ display: inline-block; direction: ltr; unicode-bidi: embed; word-break: break-all; font-family: Arial, sans-serif; }}
  </style>
</head>
<body>
  <header>
    <div>
      <h1>مصادر ومنهجية البيانات</h1>
      <p class="sub">{html.escape(REPORT_TITLE)} | {html.escape(SNAPSHOT_LABEL)}</p>
    </div>
    <img src="../assets/alforaij_logo.png" alt="{html.escape(CLIENT_NAME)}">
  </header>
  <section class="kpis">
    <div class="box"><span>حركة الدلال</span><strong>{fmt_int(metrics["total_core"])}</strong></div>
    <div class="box"><span>عروض البيع</span><strong>{fmt_int(metrics["sale_count"])}</strong></div>
    <div class="box"><span>عروض الإيجار</span><strong>{fmt_int(metrics["rent_count"])}</strong></div>
    <div class="box"><span>طلبات الشراء</span><strong>{fmt_int(metrics["buy_count"])}</strong></div>
  </section>
  <h2>منهجية مختصرة</h2>
  <ul>
    <li>هذا تقرير تشغيلي من بيانات منصة الفريج، وليس تقريراً وطنياً عن سوق العقار الكويتي.</li>
    <li>تم السحب من واجهة البحث العامة لمنصة الفريج بتاريخ {html.escape(SNAPSHOT_LABEL)}.</li>
    <li>تمت محاولة فحص تفاصيل كل إعلان لاستخراج التصنيف والمساحة والملامح المنشورة، واستخدمت التفاصيل المتاحة دون عرض بيانات شخصية.</li>
    <li>تمت قراءة أنواع المعاملات كما هي منشورة: للبيع، للإيجار، مطلوب للشراء، مطلوب للإيجار، للبدل.</li>
    <li>المتوسطات والوسيط مبنية على الأسعار المعلنة فقط؛ الأسعار الصفرية أو الفارغة تعامل كسعر غير معلن.</li>
    <li>تم استبعاد أرقام الهواتف وأسماء الملاك أو العملاء، ولم يتم عرض إلا بيانات مجمعة أو سجلات منظفة.</li>
  </ul>
  <h2>روابط المصادر</h2>
  <p>توفرت تفاصيل مباشرة لعدد {fmt_int(metrics["details_loaded"])} إعلان/طلب من أصل {fmt_int(metrics["total_all"])} سجل منشور، وبقيت السجلات الأخرى محتسبة من نتائج البحث العامة.</p>
  <table>
    <thead><tr><th>المصدر</th><th>الرابط</th></tr></thead>
    <tbody>{link_html}</tbody>
  </table>
</body>
</html>
"""
    SOURCES_HTML_PATH.write_text(html_doc, encoding="utf-8")


def find_edge() -> str | None:
    candidates = [
        shutil.which("msedge"),
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return candidate
    return None


def render_html_to_pdf(html_path: Path, pdf_path: Path) -> None:
    edge = find_edge()
    if not edge:
        raise RuntimeError("Microsoft Edge was not found; cannot render PDF.")
    user_data_dir = tempfile.mkdtemp(prefix="edge-director-manager-")
    cmd = [
        edge,
        "--headless",
        "--disable-gpu",
        "--disable-extensions",
        "--no-first-run",
        f"--user-data-dir={user_data_dir}",
        "--print-to-pdf-no-header",
        "--no-pdf-header-footer",
        f"--print-to-pdf={pdf_path}",
        html_path.resolve().as_uri(),
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=90)
    if not pdf_path.exists() or pdf_path.stat().st_size < 1000:
        raise RuntimeError(f"PDF rendering failed for {html_path.name}.")


def render_pdf_pages(pdf_path: Path) -> dict:
    doc = fitz.open(pdf_path)
    pages: list[Path] = []
    nonblank: list[bool] = []
    for idx, page in enumerate(doc, start=1):
        pix = page.get_pixmap(matrix=fitz.Matrix(0.75, 0.75), alpha=False)
        out = QA_DIR / f"page_{idx:02d}.png"
        pix.save(out)
        img = Image.open(out).convert("RGB")
        extrema = img.getextrema()
        nonblank.append(any(lo != hi for lo, hi in extrema))
        pages.append(out)

    thumb_w = 430
    thumbs = []
    for path in pages:
        img = Image.open(path).convert("RGB")
        ratio = thumb_w / img.width
        thumbs.append(img.resize((thumb_w, int(img.height * ratio)), Image.Resampling.LANCZOS))
    cols = 2
    margin = 24
    label_h = 28
    rows = math.ceil(len(thumbs) / cols)
    thumb_h = max(img.height for img in thumbs)
    montage = Image.new("RGB", (cols * thumb_w + (cols + 1) * margin, rows * (thumb_h + label_h) + (rows + 1) * margin), "white")
    draw = ImageDraw.Draw(montage)
    for idx, thumb in enumerate(thumbs):
        col = idx % cols
        row = idx // cols
        x = margin + col * (thumb_w + margin)
        y = margin + row * (thumb_h + label_h + margin)
        draw.text((x, y), f"Page {idx + 1:02d}", fill=(15, 23, 42))
        montage.paste(thumb, (x, y + label_h))
    montage.save(MONTAGE_PATH)
    return {"page_count": len(doc), "all_nonblank": all(nonblank), "pages": [str(p) for p in pages]}


def validate(metrics: dict) -> dict:
    required = [REPORT_PDF_PATH, PPTX_PATH, EXCEL_PATH, SOURCES_PDF_PATH]
    missing = [str(path) for path in required if not path.exists() or path.stat().st_size == 0]
    if missing:
        raise RuntimeError("Missing output files: " + ", ".join(missing))

    prs = Presentation(PPTX_PATH)
    ppt_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    checks = {
        "عدد شرائح العرض": len(prs.slides) == 8,
        "رقم حركة الدلال موجود": fmt_int(metrics["total_core"]) in ppt_text,
        "رقم عروض البيع موجود": fmt_int(metrics["sale_count"]) in ppt_text,
        "رقم عروض الإيجار موجود": fmt_int(metrics["rent_count"]) in ppt_text,
        "رقم طلبات الشراء موجود": fmt_int(metrics["buy_count"]) in ppt_text,
    }
    if not all(checks.values()):
        failed = [name for name, ok in checks.items() if not ok]
        raise RuntimeError("Validation failed: " + ", ".join(failed))
    return checks


def create_qa_md(metrics: dict, pdf_info: dict, checks: dict) -> None:
    lines = [
        "# QA",
        "",
        f"- تاريخ التوليد: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}",
        f"- صفحات PDF: {pdf_info['page_count']}",
        f"- فحص الصفحات غير الفارغة: {'ناجح' if pdf_info['all_nonblank'] else 'يحتاج مراجعة'}",
        f"- Montage بصري: {MONTAGE_PATH.name}",
        "",
        "## فحوص الأرقام",
    ]
    for name, ok in checks.items():
        lines.append(f"- {name}: {'ناجح' if ok else 'يحتاج مراجعة'}")
    lines.extend(
        [
            "",
            "## نطاق البيانات",
            f"- حركة الدلال ضمن نطاق الطلب: {metrics['total_core']}",
            f"- عروض البيع: {metrics['sale_count']}",
            f"- طلبات الشراء: {metrics['buy_count']}",
            f"- عروض الإيجار: {metrics['rent_count']}",
            f"- طلبات الإيجار: {metrics['rent_request_count']}",
            "",
            "## حماية البيانات",
            "- لم يتم إدراج هواتف أو أسماء ملاك/عملاء في التقرير أو سجل البيانات المنظف.",
            "- كل المتوسطات السعرية مبنية على القيم المعلنة فقط.",
        ]
    )
    QA_MD_PATH.write_text("\n".join(lines) + "\n", encoding="utf-8")


def copy_final_files() -> None:
    copies = [
        (REPORT_PDF_PATH, FINAL_REPORT_PDF),
        (PPTX_PATH, FINAL_PPTX),
        (EXCEL_PATH, FINAL_EXCEL),
        (SOURCES_PDF_PATH, FINAL_SOURCES_PDF),
    ]
    for src, dst in copies:
        shutil.copy2(src, dst)


def create_manifest(metrics: dict, pdf_info: dict) -> None:
    payload = {
        "name": "director_manager_operational_statistics",
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "snapshot_date": SNAPSHOT_DATE.isoformat(),
        "sources": {
            "front": ALFORAIJ_FRONT_URL,
            "search": ALFORAIJ_SEARCH_URL,
            "transaction_types": TRANSACTION_TYPES_URL,
            "listings": LISTINGS_URL,
        },
        "metrics": {
            "total_core": metrics["total_core"],
            "sale_count": metrics["sale_count"],
            "buy_count": metrics["buy_count"],
            "house_sale_count": metrics["house_sale_count"],
            "rent_count": metrics["rent_count"],
            "rent_request_count": metrics["rent_request_count"],
            "exchange_count": metrics["exchange_count"],
        },
        "pdf_pages": pdf_info["page_count"],
        "files_to_manager": [str(path.relative_to(OUTPUT_DIR)) for path in [FINAL_REPORT_PDF, FINAL_PPTX, FINAL_EXCEL, FINAL_SOURCES_PDF]],
    }
    MANIFEST_PATH.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def main() -> None:
    ensure_dirs()
    type_names, listings = fetch_all_data()
    metrics = build_analysis(type_names, listings)
    build_report_html(metrics)
    create_sources_html(metrics)
    render_html_to_pdf(REPORT_HTML_PATH, REPORT_PDF_PATH)
    render_html_to_pdf(SOURCES_HTML_PATH, SOURCES_PDF_PATH)
    create_pptx(metrics)
    create_excel(metrics)
    pdf_info = render_pdf_pages(REPORT_PDF_PATH)
    checks = validate(metrics)
    create_qa_md(metrics, pdf_info, checks)
    copy_final_files()
    create_manifest(metrics, pdf_info)
    print(json.dumps({"output": str(TO_MANAGER_DIR), "metrics": MANIFEST_PATH.read_text(encoding="utf-8")}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
